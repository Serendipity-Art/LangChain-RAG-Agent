"""
LangChain RAG Agent — HuggingFace Spaces Gradio UI
支持 PDF / Word / PPT / TXT / Excel / CSV 多格式
纯 RAG 模式 + Agent 智能体模式（带第三方工具）
"""

import os, time, traceback
from pathlib import Path

import gradio as gr
import pandas as pd

from langchain_core.documents import Document
from langchain_core.prompts import PromptTemplate
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai import ChatOpenAI

# ── 厂商配置 ──

PROVIDERS = {
    "DeepSeek": {
        "base_url": "https://api.deepseek.com",
        "models": ["deepseek-v4-flash", "deepseek-v4-pro"],
        "default_model": "deepseek-v4-flash",
    },
    "通义千问 (Qwen)": {
        "base_url": "https://dashscope.aliyuncs.com/compatible-mode/v1",
        "models": ["qwen-max", "qwen-plus", "qwen-turbo", "qwen-long"],
        "default_model": "qwen-max",
    },
    "智谱AI (GLM)": {
        "base_url": "https://open.bigmodel.cn/api/paas/v4",
        "models": ["glm-5.1", "glm-5", "glm-4-plus", "glm-4-air"],
        "default_model": "glm-5.1",
    },
    "月之暗面 (Kimi)": {
        "base_url": "https://api.moonshot.cn/v1",
        "models": ["kimi-k2.6", "kimi-k2.5", "moonshot-v1"],
        "default_model": "kimi-k2.6",
    },
    "MiniMax": {
        "base_url": "https://api.minimax.chat/v1",
        "models": ["MiniMax-Text-01", "abab6.5s-chat", "abab5.5s-chat"],
        "default_model": "MiniMax-Text-01",
    },
    "豆包 (火山引擎)": {
        "base_url": "https://ark.cn-beijing.volces.com/api/v3",
        "models": ["doubao-pro-128k", "doubao-pro-32k", "doubao-lite-32k"],
        "default_model": "doubao-pro-128k",
    },
}

EMBEDDING_MODELS = {
    "all-MiniLM-L6-v2": "通用（384维，80MB，速度最快，中英文均可）",
    "all-mpnet-base-v2": "通用高精度（768维，420MB，精度最高）",
    "BAAI/bge-small-zh-v1.5": "🇨🇳 中文优化·轻量版（384维，快速，中文检索优秀）",
    "BAAI/bge-base-zh-v1.5": "🇨🇳 中文优化·高精度版（768维，中文检索SOTA，推荐中文场景）",
}

# ── 工具函数 ──

def build_llm(api_key, base_url, model_name, temperature):
    if not api_key:
        return None
    return ChatOpenAI(model=model_name, openai_api_key=api_key, openai_api_base=base_url, temperature=temperature)

def build_embeddings(model_name):
    return HuggingFaceEmbeddings(model_name=model_name)

def get_filepath(file_obj):
    """统一提取文件路径"""
    if file_obj is None:
        return None
    return file_obj if isinstance(file_obj, str) else getattr(file_obj, "name", None)

def load_documents(filepath):
    """通用文档加载器：自动识别后缀，返回 Document 列表"""
    suffix = Path(filepath).suffix.lower()
    source = Path(filepath).name

    if suffix == ".txt":
        with open(filepath, "r", encoding="utf-8") as f:
            return [Document(page_content=f.read(), metadata={"source": source})]

    elif suffix == ".docx":
        import docx
        doc = docx.Document(filepath)
        text = "\n".join(p.text for p in doc.paragraphs)
        return [Document(page_content=text, metadata={"source": source})]

    elif suffix == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        return [Document(page_content=text, metadata={"source": source})]

    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(filepath)
        parts = []
        for slide in prs.slides:
            texts = [shp.text for shp in slide.shapes if hasattr(shp, "text") and shp.text.strip()]
            if texts:
                parts.append("\n".join(texts))
        return [Document(page_content="\n---\n".join(parts), metadata={"source": source})]

    # Excel / CSV → 每行转一个 Document
    elif suffix in (".xlsx", ".xls"):
        df = pd.read_excel(filepath, engine="openpyxl")
        return [
            Document(page_content=" | ".join(f"{c}: {r[c]}" for c in df.columns), metadata={"source": source})
            for _, r in df.iterrows()
        ]
    elif suffix == ".csv":
        try:
            df = pd.read_csv(filepath, encoding="utf-8")
        except UnicodeDecodeError:
            df = pd.read_csv(filepath, encoding="gbk")
        return [
            Document(page_content=" | ".join(f"{c}: {r[c]}" for c in df.columns), metadata={"source": source})
            for _, r in df.iterrows()
        ]
    else:
        raise ValueError(f"不支持的文件格式: {suffix}")

def build_vector_store(filepath, embedding_model, chunk_size=500, chunk_overlap=50):
    """统一入口：加载 → 分割 → FAISS"""
    docs = load_documents(filepath)
    if not docs:
        raise ValueError("文档内容为空")
    chunks = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap).split_documents(docs)
    if not chunks:
        raise ValueError("分割后没有产生任何文档片段")
    vs = FAISS.from_documents(chunks, embedding_model)
    return vs, len(docs), len(chunks)

# ── RAG 问答核心 ──

RAG_PROMPT = """你是一个回答机器人。根据已知信息回答用户问题。

严格要求：
1. 完全基于已知信息，不编造答案。
2. 已知信息不足时，回复"我无法回答你的问题"。
3. 不要使用 Markdown 格式，用纯文本回复。

已知信息：
{info}

用户问：{question}

请用中文回答："""

def answer_rag(query, vector_store, llm, k=3):
    t0 = time.time()
    docs = vector_store.as_retriever(search_kwargs={"k": k}).invoke(query)
    t1 = time.time()
    ans = llm.invoke(PromptTemplate.from_template(RAG_PROMPT).format(
        info="\n\n".join(d.page_content for d in docs), question=query)).content
    t2 = time.time()
    return round(t1-t0,3), round(t2-t1,3), round(t2-t0,3), ans, docs

def answer_agent(query, vector_store, llm, k=3):
    from langchain.agents import create_agent
    from langchain.tools import tool
    import datetime, math

    @tool
    def search_knowledge(q: str) -> str:
        """📖 搜索本地知识库：从上传文档中找到最相关的信息"""
        docs = vector_store.as_retriever(search_kwargs={"k": k}).invoke(q)
        return "\n\n".join(d.page_content for d in docs) if docs else "未找到相关信息"

    @tool
    def calculator(expression: str) -> str:
        """🧮 数学计算：输入数学表达式，返回计算结果。例：'123*456'"""
        try:
            return str(eval(expression, {"__builtins__": {}}, {"math": math}))
        except Exception as e:
            return f"计算失败: {e}"

    @tool
    def get_current_time() -> str:
        """🕐 获取当前日期和时间"""
        return datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    agent = create_agent(
        model=llm,
        tools=[search_knowledge, calculator, get_current_time],
        system_prompt=(
            "你是智能助手，可调用多个工具回答用户问题。\n"
            "1. 知识类问题 → 搜索知识库\n"
            "2. 数学问题 → 计算器\n"
            "3. 时间问题 → 获取当前时间\n"
            "若搜索知识库没有足够信息，请说'我无法回答'。")
    )
    result = agent.invoke({"messages": [{"role": "user", "content": query}]})
    last = result["messages"][-1]
    return last.content if hasattr(last, "content") else str(last)

# ── UI 处理函数 ──

def process_upload(file_obj, chunk_size, chunk_overlap, embedding_name, progress=gr.Progress()):
    if file_obj is None:
        return None, '<div class="status-msg warn">⚠️ 请先选择一个文件</div>'
    fp = get_filepath(file_obj)
    if not fp:
        return None, '<div class="status-msg error">❌ 无法读取文件</div>'

    suffix = Path(fp).suffix.lower()
    SUPPORTED = {".txt", ".docx", ".pdf", ".pptx", ".xlsx", ".xls", ".csv"}
    if suffix not in SUPPORTED:
        return None, f'<div class="status-msg error">❌ 不支持 {suffix}，请上传 TXT/DOCX/PDF/PPTX/Excel/CSV</div>'

    progress(0.1, desc="加载文档...")
    emb = build_embeddings(embedding_name)
    try:
        vs, n_docs, n_chunks = build_vector_store(fp, emb, chunk_size, chunk_overlap)
    except Exception as e:
        return None, f'<div class="status-msg error">❌ 处理失败：{e}</div>'

    fname = Path(fp).name
    return vs, f'<div class="status-msg success">✅ {fname} → {n_chunks} 个片段，向量库就绪</div>'

def answer_question(query, vs, api_key, base_url, model_name, temperature, k, mode):
    if vs is None:
        return "⚠️ 请先上传文档构建向量库", "", "", ""
    if not api_key:
        return "⚠️ 请设置 API Key", "", "", ""
    llm = build_llm(api_key, base_url, model_name, temperature)
    if llm is None:
        return "⚠️ LLM 初始化失败", "", "", ""
    try:
        if mode == "Agent 智能体":
            t0 = time.time()
            ans = answer_agent(query, vs, llm, k)
            el = round(time.time()-t0, 3)
            s = f'<div class="stat-grid"><div class="stat-item"><span class="stat-label">⏱️ 总耗时</span><span class="stat-value">{el}s</span></div><div class="stat-item"><span class="stat-label">🤖 模式</span><span class="stat-value">Agent</span></div></div>'
            return ans, s, "Agent 模式下 LLM 自动调度多个工具（知识库搜索/计算器/时间）", ""
        rt, lt, tt, ans, docs = answer_rag(query, vs, llm, k)
        s = (f'<div class="stat-grid">'
             f'<div class="stat-item"><span class="stat-label">🔍 检索</span><span class="stat-value">{rt}s</span></div>'
             f'<div class="stat-item"><span class="stat-label">🧠 LLM</span><span class="stat-value">{lt}s</span></div>'
             f'<div class="stat-item"><span class="stat-label">⏱️ 总计</span><span class="stat-value">{tt}s</span></div>'
             f'<div class="stat-item"><span class="stat-label">📄 片段</span><span class="stat-value">{len(docs)}</span></div></div>')
        d = "---\n\n".join(f"**片段 {i+1}**\n{d.page_content[:200]}..." for i,d in enumerate(docs))
        return ans, s, d, ""
    except Exception as e:
        return f"❌ 回答失败：{traceback.format_exc()}", "", "", ""

def update_provider(name):
    info = PROVIDERS.get(name, PROVIDERS["DeepSeek"])
    return gr.update(value=info["base_url"]), gr.update(choices=info["models"], value=info["default_model"])

# ── 构建 UI ──

def build_ui():
    names = list(PROVIDERS.keys())
    default = PROVIDERS[names[0]]

    with gr.Blocks(title="LangChain RAG Agent") as demo:
        # ── 顶部标题 ──
        gr.HTML("""
        <div class="header-banner">
            <div class="header-inner">
                <div class="header-icon">🦜</div>
                <div class="header-text">
                    <h1 class="header-title">LangChain RAG Agent</h1>
                    <p class="header-sub">上传文档 → 自动建库 → 基于知识库的智能问答</p>
                </div>
            </div>
            <div class="header-tags">
                <span>TXT</span><span>Word</span><span>PDF</span>
                <span>PPT</span><span>Excel</span><span>CSV</span>
            </div>
        </div>
        """)

        vs_state = gr.State(None)

        with gr.Row(equal_height=False):
            # ═══ 左栏 ═══
            with gr.Column(scale=35, min_width=300, elem_classes="col-card"):
                gr.HTML('<div class="col-title"><span>⚙️</span> 配置面板</div>')

                gr.HTML('<div class="section-header" style="border-color:#6366f1"><span>🔑</span> API 设置</div>')
                provider_dd = gr.Dropdown(label="大模型厂商", choices=names, value=names[0])
                base_url = gr.Textbox(label="Base URL", value=default["base_url"], interactive=False)
                model_dd = gr.Dropdown(label="模型", choices=default["models"], value=default["default_model"])
                provider_dd.change(fn=update_provider, inputs=provider_dd, outputs=[base_url, model_dd])
                api_key = gr.Textbox(label="API Key", type="password", placeholder="填入对应厂商的 API Key", value=os.getenv("OPENAI_API_KEY1", ""))
                temperature = gr.Slider(label="Temperature（温度）", minimum=0, maximum=2, value=0.7, step=0.1)

                gr.HTML('<div class="section-divider"></div>')

                gr.HTML('<div class="section-header" style="border-color:#10b981"><span>📂</span> 上传文档</div>')
                file_input = gr.File(label="拖拽或点击上传文件", file_types=[".txt", ".docx", ".pdf", ".pptx", ".xlsx", ".xls", ".csv"])
                with gr.Accordion("⚡ 分割参数", open=False):
                    with gr.Row():
                        chunk_size = gr.Slider(label="Chunk Size", minimum=200, maximum=2000, value=500, step=50)
                        chunk_overlap = gr.Slider(label="Overlap", minimum=0, maximum=400, value=50, step=10)
                emb_dd = gr.Dropdown(
                    label="🧠 Embedding 模型",
                    choices=[f"{k} — {v}" for k,v in EMBEDDING_MODELS.items()],
                    value="all-MiniLM-L6-v2 — 通用（384维，80MB，速度最快，中英文均可）")
                gr.HTML('<div class="tip-text">💡 中文场景推荐: <code>BAAI/bge-base-zh-v1.5</code></div>')
                build_btn = gr.Button("🚀 构建向量库", variant="primary")
                build_status = gr.HTML('<div class="status-msg idle">💡 上传文件后点击构建</div>')

                gr.HTML('<div class="section-divider"></div>')

                gr.HTML('<div class="section-header" style="border-color:#f59e0b"><span>🔍</span> 回答模式</div>')
                k_slider = gr.Slider(label="检索 Top-K", minimum=1, maximum=10, value=3, step=1)
                mode_radio = gr.Radio(
                    label="模式",
                    choices=["纯 RAG（检索+生成）", "Agent 智能体"],
                    value="纯 RAG（检索+生成）",
                )
                gr.HTML('<div class="tip-text" style="margin-top:4px">💡 Agent 模式可调用知识库搜索 + 计算器 + 时间查询</div>')

            # ═══ 右栏 ═══
            with gr.Column(scale=65, elem_classes="col-card"):
                gr.HTML('<div class="col-title"><span>💬</span> 对话</div>')

                chatbot = gr.Chatbot(label="", height=460, avatar_images=[None, None])
                query_input = gr.Textbox(label="", placeholder="输入你的问题，例如：苏州园林有什么特点？", lines=2)

                with gr.Row():
                    submit_btn = gr.Button("🎯 发送", variant="primary", scale=2)
                    clear_btn = gr.Button("🔄 清空", scale=1)

                with gr.Accordion("📊 性能统计", open=True):
                    stats_output = gr.HTML('<div class="stat-placeholder">等待首次问答...</div>')
                with gr.Accordion("📄 检索溯源", open=False):
                    docs_display = gr.Markdown("展开查看检索原文")

                gr.HTML('<div class="examples-title">💡 试试这些问题</div>')
                gr.Examples(
                    examples=[["苏州园林有什么特点？"], ["虎丘塔的历史背景是什么？"],
                              ["苏州园林和北方园林有什么区别？"], ["拙政园是谁建造的？"]],
                    inputs=query_input)

                gr.HTML("""
                <div class="footer-card">
                    <div class="footer-card-title">🚀 项目概述</div>
                    <div class="footer-desc"><b>LangChain RAG Agent</b> 开箱即用的检索增强生成系统。支持 TXT / Word / PDF / PPT / Excel / CSV 六种格式，上传文档自动建库，即问即答。</div>
                    <div class="footer-card-title" style="margin-top:22px">🎯 应用场景</div>
                    <div class="footer-apps">
                        <div class="footer-app-item"><div class="footer-app-icon">📚</div><div class="footer-app-text"><b>学术研究</b><span>论文问答、文献综述</span></div></div>
                        <div class="footer-app-item"><div class="footer-app-icon">🏫</div><div class="footer-app-text"><b>课堂教学</b><span>课件答疑、资源管理</span></div></div>
                        <div class="footer-app-item"><div class="footer-app-icon">🏢</div><div class="footer-app-text"><b>企业培训</b><span>知识库、员工问答</span></div></div>
                        <div class="footer-app-item"><div class="footer-app-icon">📊</div><div class="footer-app-text"><b>数据分析</b><span>报表解读、数据查询</span></div></div>
                        <div class="footer-app-item"><div class="footer-app-icon">📝</div><div class="footer-app-text"><b>文档管理</b><span>合同检索、政策问答</span></div></div>
                        <div class="footer-app-item"><div class="footer-app-icon">⚖️</div><div class="footer-app-text"><b>法律法务</b><span>法规查询、案例检索</span></div></div>
                    </div>
                    <div class="footer-card-title" style="margin-top:22px">📌 技术路线图</div>
                    <div class="roadmap-list">
                        <div class="roadmap-item"><span class="roadmap-dot" style="background:#6366f1"></span><span class="roadmap-text"><b>Phase 1 ✅ 已实现</b> —— 6 种格式支持、纯 RAG / Agent 双模式、中文 Embedding 优化</span></div>
                        <div class="roadmap-item"><span class="roadmap-dot" style="background:#10b981"></span><span class="roadmap-text"><b>Phase 2 · 体验优化</b> —— 流式输出（感知 3s→0.3s）、多轮对话记忆、FAISS 持久化</span></div>
                        <div class="roadmap-item"><span class="roadmap-dot" style="background:#f59e0b"></span><span class="roadmap-text"><b>Phase 3 · 性能飞跃</b> —— 🔥 Redis 缓存（热门问答毫秒返回，API 成本降 60%）、混合检索（BM25+向量，召回率 +30%）</span></div>
                        <div class="roadmap-item"><span class="roadmap-dot" style="background:#8b5cf6"></span><span class="roadmap-text"><b>Phase 4 · 架构升级</b> —— 分布式向量库（百万级语料）、多用户隔离、Web 网页抓取</span></div>
                    </div>
                    <div class="footer-contact"><span class="contact-icon">📧</span><span class="contact-text">zhujunheng8@qq.com</span><span class="contact-divider">|</span><span class="contact-year">2026 LangChain RAG Agent</span></div>
                </div>
                """)

        # ── 事件 ──
        build_btn.click(fn=process_upload,
            inputs=[file_input, chunk_size, chunk_overlap, emb_dd],
            outputs=[vs_state, build_status]).then(fn=lambda: [], outputs=[chatbot])

        def respond(query, history, vs, ak, bu, mn, te, k, md):
            if not query or not query.strip():
                return history, "", ""
            ans, st, docs, err = answer_question(query, vs, ak, bu, mn, te, k, md)
            if err and err.startswith("⚠️"):
                history = history or []
                history.append({"role":"user","content":[{"type":"text","text":query}]})
                history.append({"role":"assistant","content":[{"type":"text","text":f"❌ {err}"}]})
                return history, "", ""
            history = history or []
            history.append({"role":"user","content":[{"type":"text","text":query}]})
            history.append({"role":"assistant","content":[{"type":"text","text":ans}]})
            return history, st, docs

        for trigger in [submit_btn.click, query_input.submit]:
            trigger(fn=respond,
                inputs=[query_input, chatbot, vs_state, api_key, base_url, model_dd, temperature, k_slider, mode_radio],
                outputs=[chatbot, stats_output, docs_display]).then(fn=lambda:"", outputs=[query_input])

        clear_btn.click(fn=lambda: ([], '<div class="stat-placeholder">等待首次问答...</div>', ""),
                        outputs=[chatbot, stats_output, docs_display])

    return demo

# ── 启动 ──

if __name__ == "__main__":
    demo = build_ui()
    port = int(os.getenv("GRADIO_SERVER_PORT", 7860))
    demo.launch(server_port=port, server_name="0.0.0.0", theme=gr.themes.Soft(),
        css="""
/* ── 全局 ── */
.gradio-container { max-width:1440px !important; margin:0 auto !important; padding:24px 32px !important;
    background:linear-gradient(135deg,#eef2ff 0%,#fdf2f8 40%,#ede9fe 75%,#e0f2fe 100%) !important; }
footer { display:none !important; }

/* ── 标题栏 ── */
.header-banner {
    background:linear-gradient(135deg,#1e1b4b 0%,#312e81 50%,#3730a3 100%);
    border-radius:18px; padding:28px 32px; margin-bottom:24px;
    box-shadow:0 8px 32px rgba(79,70,229,0.15);
}
.header-inner { display:flex; align-items:center; gap:18px; }
.header-icon { font-size:44px; width:68px; height:68px; background:rgba(255,255,255,0.1);
    border-radius:16px; display:flex; align-items:center; justify-content:center; flex-shrink:0; }
.header-title { margin:0; color:#fff !important; font-size:26px; font-weight:700; letter-spacing:0.5px; }
.header-sub { margin:4px 0 0; color:#fff !important; font-size:14px; font-weight:500; }
.header-tags { margin-top:16px; display:flex; gap:8px; flex-wrap:wrap; }
.header-tags span { background:rgba(255,255,255,0.15); color:#fff !important;
    padding:4px 14px; border-radius:20px; font-size:12px; border:1px solid rgba(255,255,255,0.15); }

/* ── 卡片 ── */
.col-card > div { background:#fff !important; border-radius:16px !important;
    padding:8px 20px 20px !important;
    box-shadow:0 1px 4px rgba(0,0,0,0.04),0 4px 16px rgba(0,0,0,0.03) !important;
    border:1px solid rgba(0,0,0,0.04) !important; height:fit-content !important; }
.col-title { display:flex; align-items:center; gap:8px; font-size:17px; font-weight:700; color:#1f2937; padding:12px 0 4px; }
.section-header { display:flex; align-items:center; gap:8px; font-size:14px; font-weight:600; color:#374151;
    border-left:3px solid #6366f1; padding:8px 0 8px 14px; margin:14px 0 8px; }
.section-divider { height:1px; background:linear-gradient(to right,#e5e7eb,transparent); margin:18px 0; }
.col-card .block { box-shadow:none !important; background:transparent !important; }
input,textarea,select { border-radius:8px !important; }

/* ── 按钮 ── */
button.primary { border-radius:10px !important; font-weight:600 !important; padding:10px 24px !important;
    background:linear-gradient(135deg,#6366f1,#4f46e5) !important; border:none !important; transition:all .2s !important; }
button.primary:hover { transform:translateY(-1px) !important;
    box-shadow:0 6px 20px rgba(79,70,229,0.35) !important; }

/* ── 聊天 ── */
.chatbot { border-radius:12px !important; border:1px solid #e5e7eb !important; }

/* ── 状态 ── */
.status-msg { padding:10px 14px; border-radius:10px; font-size:13px; margin:4px 0; }
.status-msg.idle { background:#f3f4f6; color:#6b7280; }
.status-msg.success { background:#ecfdf5; color:#065f46; border:1px solid #a7f3d0; }
.status-msg.warn { background:#fffbeb; color:#92400e; border:1px solid #fde68a; }
.status-msg.error { background:#fef2f2; color:#991b1b; border:1px solid #fecaca; }

/* ── 统计 ── */
.stat-grid { display:flex; gap:10px; flex-wrap:wrap; }
.stat-item { background:#f8fafc; border:1px solid #e2e8f0; border-radius:10px;
    padding:10px 16px; display:flex; flex-direction:column; align-items:center; min-width:80px; flex:1; }
.stat-label { font-size:11px; color:#94a3b8; text-transform:uppercase; letter-spacing:0.5px; }
.stat-value { font-size:18px; font-weight:700; color:#1e293b; margin-top:2px; }
.stat-placeholder { padding:10px 14px; background:#f9fafb; border-radius:8px; font-size:13px; color:#9ca3af; }

/* ── 提示 ── */
.tip-text { font-size:12px; color:#888; padding:0 4px; margin-bottom:8px; }
.tip-text code { background:#ecfdf5; padding:1px 6px; border-radius:4px; font-size:11px; }

/* ── 示例 ── */
.examples-title { font-weight:600; font-size:14px; color:#1f2937; margin-top:12px; margin-bottom:4px; }
.examples-container button { border-radius:20px !important; border:1px solid #e5e7eb !important;
    background:#fff !important; transition:all .15s !important; }
.examples-container button:hover { border-color:#6366f1 !important; color:#6366f1 !important;
    background:#f5f3ff !important; }

/* ── 底部卡片 ── */
.footer-card { margin-top:28px; padding:36px 36px;
    background:linear-gradient(135deg,#f8fafc 0%,#f1f5f9 100%);
    border:1px solid #e2e8f0; border-radius:18px; }
.footer-card-title { font-size:20px; font-weight:700; color:#1e293b; margin-bottom:16px; }
.footer-desc { font-size:15px; color:#475569; line-height:1.9; margin-bottom:6px; }
.footer-apps { display:grid; grid-template-columns:1fr 1fr; gap:10px; margin-bottom:6px; }
.footer-app-item { display:flex; align-items:center; gap:12px; padding:14px 18px;
    background:#fff; border-radius:12px; border:1px solid #e5e7eb; transition:all .15s; }
.footer-app-item:hover { border-color:#6366f1; box-shadow:0 2px 12px rgba(99,102,241,0.1); }
.footer-app-icon { font-size:24px; flex-shrink:0; }
.footer-app-text { display:flex; flex-direction:column; font-size:15px; color:#374151; }
.footer-app-text span { font-size:13px; color:#94a3b8; margin-top:2px; }
.roadmap-list { display:flex; flex-direction:column; gap:8px; margin-bottom:22px; }
.roadmap-item { display:flex; align-items:flex-start; gap:12px; padding:12px 18px;
    background:#fff; border-radius:12px; border:1px solid #e5e7eb; }
.roadmap-dot { width:10px; height:10px; border-radius:50%; flex-shrink:0; margin-top:7px; }
.roadmap-text { font-size:14px; color:#475569; line-height:1.7; }
.footer-contact { display:flex; align-items:center; gap:12px; padding:14px 22px;
    background:linear-gradient(135deg,#eef2ff,#f5f3ff); border-radius:12px;
    border:1px solid #ddd6fe; font-size:15px; color:#4338ca; }
.contact-icon { font-size:20px; }
.contact-divider { color:#c7d2fe; font-size:16px; }
.contact-year { color:#94a3b8; font-size:13px; }
""")
